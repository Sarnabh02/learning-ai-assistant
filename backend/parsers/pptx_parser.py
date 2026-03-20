import io
from pptx import Presentation


def parse_pptx_bytes(data: bytes) -> tuple[str | None, str | None]:
    """Extract text from PPTX bytes. Returns (text, error)."""
    try:
        prs = Presentation(io.BytesIO(data))
        slide_texts: list[str] = []
        for slide in prs.slides:
            shapes_text: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes_text.append(shape.text.strip())
            if shapes_text:
                slide_texts.append("\n".join(shapes_text))
        text = "\n\n".join(slide_texts)
        if not text.strip():
            return None, (
                "No text found in this presentation. "
                "It may contain only images or diagrams."
            )
        return text, None
    except Exception as e:
        return None, f"PPTX parsing failed: {e}"
